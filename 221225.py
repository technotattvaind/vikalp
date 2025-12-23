#!/usr/bin/env python
# coding: utf-8

# In[23]:


# ============================================================
# SIRF Analytics – Diploma Sector
# Premium Kesariya Dashboard | Excel-1 Aligned
# ============================================================

import os
import re
import numpy as np
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
import streamlit as st
import streamlit.components.v1 as components
import folium
from streamlit_folium import st_folium

#---add for ppt
\
from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO
import plotly.io as pio
import tempfile


from sklearn.cluster import KMeans
from sklearn.preprocessing import StandardScaler
from sklearn.decomposition import PCA
from sklearn.ensemble import RandomForestRegressor
# ============================================================
# PAGE CONFIG
# ============================================================
st.set_page_config(
    page_title="SIRF Analytics – Diploma Sector",
    layout="wide",
    initial_sidebar_state="expanded"
)

# 🔥 GLOBAL ZOOM (MOST IMPORTANT)
st.markdown("""
<script>
document.body.style.zoom = "1.05";
</script>
""", unsafe_allow_html=True)

# ============================================================
# PREMIUM KESARIYA CSS + GLOBAL FONT SCALE (FIXED)
# ============================================================


st.markdown("""
<style>

/* ===================== HEADINGS ===================== */
h1 {
    font-size: 7.2rem !important;
    font-weight: 600 !important;
    background: linear-gradient(90deg,#FF8F00,#FF6F00);
    -webkit-background-clip: text;
    -webkit-text-fill-color: transparent;
}
h2 { font-size: 2.4 rem !important; font-weight: 700; color:#BF360C; }
h3 { font-size: 1.8 rem !important; font-weight:600; color:#5D4037; }

/* =====================================================
   🔥 REAL GLOBAL FONT FIX (THIS WAS MISSING)
===================================================== */
section.main * {
    font-size: 20px !important;
    line-height: 1.65;
}

/* Markdown */
.stMarkdown p,
.stMarkdown li,
.stMarkdown span {
    font-size: 23px !important;
}

/* Captions */
.stCaption {
    font-size: 20px !important;
    color: #6D4C41;
}

/* Sidebar */
section[data-testid="stSidebar"] * {
    font-size: 19px !important;
}

/* Inputs */
div[data-baseweb="select"] span,
div[data-baseweb="input"] input,
div[data-baseweb="textarea"] textarea {
    font-size: 16px !important;
}

/* Buttons */
button {
    font-size: 20px !important;
    font-weight: 700;
}

/* =====================================================
   PREMIUM LARGE STREAMLIT TABS — FINAL FIX
===================================================== */

/* Tabs container */
div[data-testid="stTabs"] {
    margin-top: 20px;
    border-bottom: 3px solid #FFCC80;
    padding-bottom: 14px;
    min-height: 96px !important;          /* 🔥 force height */
}

/* Each tab wrapper */
div[data-baseweb="tab"] {
    padding: 0 !important;
    min-height: 88px !important;          /* 🔥 force height */
}

/* Actual clickable tab */
div[data-baseweb="tab"] > button {
    min-height: 88px !important;          /* 🔥 KEY FIX */
    padding: 22px 42px !important;
    border-radius: 20px !important;
    align-items: center !important;
}

/* Tab text */
div[data-baseweb="tab"] > button > span {
    font-size: 32px !important;
    font-weight: 800 !important;
    line-height: 1.3 !important;           /* 🔥 prevents squeeze */
    color: #6D4C41 !important;
}

/* Active tab */
div[data-baseweb="tab"][aria-selected="true"] > button {
    background: linear-gradient(90deg,#FFE0B2,#FFCC80) !important;
    box-shadow: 0 10px 22px rgba(0,0,0,0.18);
}

/* Active tab text */
div[data-baseweb="tab"][aria-selected="true"] > button > span {
    font-size: 34px !important;
    font-weight: 900 !important;
    color: #BF360C !important;
}

/* Hover */
div[data-baseweb="tab"] > button:hover {
    background-color: #FFF3E0 !important;
}

   🔠 FORCE TAB FONT SIZE (REAL TARGET)
===================================================== */

/* All tab labels */
div[data-testid="stTabs"] button[role="tab"] span {
    font-size: 36px !important;   /* 👈 increase here */
    font-weight: 900 !important;
    line-height: 1.35 !important;
}

/* Active tab label */
div[data-testid="stTabs"] button[aria-selected="true"] span {
    font-size: 38px !important;   /* 👈 bigger active tab */
    font-weight: 900 !important;
    color: #BF360C !important;
}

/* =====================================================
   METRICS & DATAFRAME (ALWAYS SMALL OTHERWISE)
===================================================== */
.stMetric label {
    font-size: 22px !important;
}

.stMetric div {
    font-size: 44px !important;
    font-weight: 900;
}

div[data-testid="stDataFrame"] {
    font-size: 18px !important;
}







/*=====================================================
   🔲 STREAMLIT SELECTBOX — FINAL STABLE FIX
===================================================== */

div[data-baseweb="select"] > div {
    min-height: 68px !important;
    background-color: #FFF8E1 !important;
}

div[data-baseweb="select"] input {
    color: #4E342E !important;
    font-size: 22px !important;
    font-weight: 700 !important;
    opacity: 1 !important;
    caret-color: #BF360C !important;
}

div[data-baseweb="select"] input:disabled {
    -webkit-text-fill-color: #4E342E !important;
    opacity: 1 !important;
}

div[data-baseweb="select"] > div > div {
    padding: 14px 18px !important;
}

div[data-baseweb="select"] svg {
    width: 26px !important;
    height: 26px !important;
}

div[data-baseweb="select"]:focus-within > div {
    border: 2px solid #FF8F00 !important;
}

/* =====================================================
   🎯 FIX: SELECTED VALUE TEXT COLOR (ONLY)
===================================================== */

div[data-baseweb="select"] input {
    color: #000000 !important;          /* pure black – always visible */
    -webkit-text-fill-color: #000000 !important;
    opacity: 1 !important;
}
</style>
""", unsafe_allow_html=True)


# ============================================================
# FILE PATHS (GLOBAL & SAFE)
# ============================================================
DATA_FILE = r"data/SIRF_DIPLOMA_SESSION_WISE_DATA10122025.xlsx"
MASTER_FILE = r"data/MASTER_INSTITUTE_LIST.xlsx"


# ============================================================
# COLUMN CONTRACT (EXCEL-1)
# ============================================================
COL_MONTH = "Month"
COL_INST = "Institute"
COL_CODE = "Inst_Code"
COL_ZONE = "Zone"

COL_TOTAL = "Total Score"
COL_TLR = "TLR Score"
COL_PPRPI = "PPRPI Score"
COL_DO = "DO Score"

COL_GRADE = "Grade"
COL_NS = "Not_Submitted"
COL_RANK = "Rank"


# ============================================================
# HUMAN-READABLE SUB-PARAMETER LABELS
# ============================================================

SUB_PARAM_LABELS = {

    # =====================================================
    # TLR — Teaching, Learning & Resources
    # =====================================================

    # ---------- Faculty Quality & Engagement ----------
    "faculty student ration": "Faculty–Student Ratio (FSR)",
    "b.tech/B.E./M.A./M.Sc": "Faculty with UG / PG Qualification",
    "m.tech/M.E./NET/SLET/SET": "Faculty with M.Tech / NET / SET",
    "P.HD": "Faculty with Ph.D.",
    "Experience < 2 Years": "Faculty Experience < 2 Years",
    "2 < = Experience <= 5 Years": "Faculty Experience 2–5 Years",
    "5 < Experience <= 10 Years": "Faculty Experience 5–10 Years",
    "Experience > 10 Years": "Faculty Experience > 10 Years",
    "MOOCs Courses,Swayam Portal, NPTEL, NITTT Modules": "Faculty MOOCs / NPTEL / NITTT Participation",
    "Departmental Training": "Departmental Training Coverage",
    "Other than departmental training": "External Training Participation",
    "Faculty pursuing Higher Education": "Faculty Pursuing Higher Education",
    "Monthly Average Students Attendance": "Average Student Attendance (%)",

    # ---------- Infrastructure: Basic Amenities ----------
    "RO Drinking Water": "RO Drinking Water Availability",
    "electricity upto 16 Hrs/Day": "Electricity Availability (≥16 Hours)",
    "Backup Availability": "Power Backup Availability",
    "Toilet in Administrative Block": "Administrative Block Toilets",
    "Toilets for Boys": "Toilets for Boys (As per Norms)",
    "Toilets for Girls": "Toilets for Girls (As per Norms)",
    "Sweeper": "Sanitation Staff Availability",
    "Toilets for Male Faculty/Staff": "Male Staff Toilets",
    "Toilets for Female Faculty/Staff": "Female Staff Toilets",
    "Internet": "Internet Availability & Speed",
    "Static IP": "Static IP Availability",
    "Wi-Fi": "Campus Wi-Fi Availability",
    "Gardening": "Green Campus / Plantation",
    "Boundary Wall": "Boundary Wall Availability",
    "CCTV Cameras": "CCTV Surveillance",
    "Guards": "Security Guards Availability",

    # ---------- Hostel Facilities ----------
    "Bathroom for Boys": "Hostel Bathrooms for Boys",
    "Bathroom for Girls": "Hostel Bathrooms for Girls",
    "Mess": "Hostel Mess Availability",
    "Seating capacity": "Mess Seating Capacity",
    "Mess Menu": "Mess Menu Displayed",
    "Student's Feedback": "Hostel Mess Student Feedback",

    # ---------- Laboratories & Workshops ----------
    "Number of Labs/Workshops required": "Required Labs / Workshops (As per Norms)",
    "Total no. of Labs/Workshops available in the institute": "Available Labs / Workshops",
    "Total no. of Tools/Machinery required": "Required Tools / Machinery",
    "Total no. of Tools/Machinery available": "Available Tools / Machinery",
    "Total no. of Tools/Machinery working": "Working Tools / Machinery",
    "Total no. of Computers required": "Required Computers",
    "Total no. of Computers available": "Available Computers",
    "Total no. of Computers working": "Working Computers",

    # ---------- Library & Accessibility ----------
    "E-Library with E-Granthalaya": "E-Library (e-Granthalaya)",
    "Books as per norms": "Library Books as per Norms",
    "Journals/Articles/Magazines": "Journals / Articles / Magazines",
    "Dispensary": "Dispensary Facility",
    "RAMP": "Ramp for Differently-Abled",
    "Accessible Toilet": "Accessible Toilets",
    "Smart Class": "Smart Classroom Availability",

    # ---------- Extra-Curricular ----------
    "Number of Events": "Extra-Curricular Events Conducted",
    "Participation at various level": "Participation at Various Levels",
    "Percentage of students participated in skill development activities": "Skill Development Participation (%)",
    "NCC/NSS/Blood Donation": "Social & Community Engagement",
    "Industry visits organized by institute": "Industry Exposure & Visits",

    # ---------- Examination & Evaluation ----------
    "Sessional Examination": "Sessional Examinations Conducted",
    "No of students assigned for examination": "Students Assigned for Examination",
    "No of Evaluators": "Number of Evaluators",

    # =====================================================
    # PPRPI — Practical, Project, Research & Impact
    # =====================================================

    # ---------- Project & Practical ----------
    "Project to faculty engagement ratio in each discipline": "Project–Faculty Engagement Ratio",
    "Total number of Jobs/Practicals Scheduled": "Total Practicals Scheduled (Monthly)",
    "Group size formed in the institute for the project work": "Project Group Size (As per Norms)",
    "Monthly Progress Report of project work/models": "Monthly Project Progress Monitoring",

    # ---------- Industry & Projects ----------
    "Institute-Industry Relationship Index": "Institute–Industry Relationship Index",
    "Project Exhibition Index": "Project Exhibition Index",
    "Student Appreciation Index": "Student Project Appreciation",
    "Institute Competitive Index": "Institute Competitiveness Index",

    # ---------- Research ----------
    "No of Journal/Book/Chapter/Conference": "Research Output Count",
    "Publications Index": "Research Publications Index",
    "Intellectual Property Right": "IPR / Patent Activities",

    # ---------- Events ----------
    "Workshop Index": "Workshops Conducted",
    "Seminar Index": "Seminars Conducted",

    # =====================================================
    # DO — Diploma Outcomes
    # =====================================================

    # ---------- Academic Outcomes ----------
    "On Time Diploma Completion Rate": "On-Time Diploma Completion Rate",
    "Percentage of students passed with honours": "Diploma with Honours (%)",
    "Total Passout Student": "Total Pass-Out Students",

    # ---------- Placement ----------
    "Companies contacted for placement": "Companies Contacted for Placement",
    "Companies visited for placement": "Companies Visiting for Placement",
    "Mock interviews for students conducted": "Mock Interviews Conducted",
    "No of Students placed through campus placement": "Campus Placements",
    "No of Students placed other than campus placement": "Off-Campus Placements",
    "Placement Success Rate": "Overall Placement Success Rate",
    "Pool Campus Status": "Pool Campus Participation",

    # ---------- Higher Education & Alumni ----------
    "No of Students opted higher education": "Students Opting for Higher Education",
    "Participation and Contribution of Alumni": "Alumni Participation & Contribution",

    # ---------- Startup & Innovation ----------
    "Has the institute established Incubation Center": "Incubation Centre Established",
    "Total number of proposals received by incubation center": "Incubation Proposals Received",
    "Has the startup established": "Startups Established",
    "Self employments through different Govt. Schemes": "Self-Employment via Govt. Schemes"
}





# ============================================================
# COLOR PALETTES (KESARIYA)
# ============================================================
KESARIYA = ["#FF9933", "#FFB84D", "#FFCC66", "#FFD699", "#FFE5CC"]

GRADE_COLORS = {
    "A": "#2E7D32",
    "B": "#F9A825",
    "C": "#FB8C00",
    "D": "#C62828",
    "NS": "#9E9E9E"
}

ZONE_COLORS_DISTINCT = {
     "WEST": "#D97706",        # Muted Saffron (authority, governance)
    "CENTRAL": "#B45309",     # Burnt Amber (core zone emphasis)
    "EAST": "#166534",        # Dark Green (stability & growth)
    "BUNDELKHAND": "#334155", # Slate Blue-Grey (development focus)
    "OTHER": "#9CA3AF"
}


# ============================================================
# HELPERS
# ============================================================
def parse_month(x):
    try:
        return pd.to_datetime(x).normalize()
    except:
        return pd.NaT

def derive_zone(code):
    if pd.isna(code):
        return "OTHER"
    return {"1":"WEST","2":"CENTRAL","3":"BUNDELKHAND","4":"EAST"}.get(str(code)[0], "OTHER")

def grade_from_score(x):
    if x >= 75: return "A"
    if x >= 50: return "B"
    if x >= 25: return "C"
    return "D"

def df_index(df):
    df = df.copy()
    if "S.No" not in df.columns:
        df.insert(0, "S.No", range(1, len(df) + 1))
    return df.reset_index(drop=True)

#for ppt 

def save_fig(fig, name):
    path = os.path.join(tempfile.gettempdir(), f"{name}.png")
    pio.write_image(fig, path, width=1400, height=800, scale=2)
    return path


from pptx.util import Pt

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

    # Improve title font
    for run in slide.shapes.title.text_frame.paragraphs[0].runs:
        run.font.size = Pt(36)
        run.font.bold = True

    return slide
#----------------------------------------------

from pptx.util import Pt

def add_text_slide(prs, title_text, body_text):
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    slide.shapes.title.text = title_text
    body = slide.placeholders[1]
    body.text = body_text

    # Professional font sizing
    for para in body.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(20)

    return slide

#-------------------------------------

from pptx.util import Inches

def add_chart_slide(prs, title, fig, img_name):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title

    img = save_fig(fig, img_name)
    slide.shapes.add_picture(
        img,
        Inches(0.8),
        Inches(1.4),
        width=Inches(8.5)
    )

    return slide


#------------------------------------------


from pptx.util import Inches, Pt

def add_table_slide(prs, title, df, max_rows=20):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title

    df = df.head(max_rows)
    rows, cols = df.shape

    table = slide.shapes.add_table(
        rows + 1,
        cols,
        Inches(0.4),
        Inches(1.5),
        Inches(9.2),
        Inches(4.8)
    ).table

    # Header
    for c, col in enumerate(df.columns):
        cell = table.cell(0, c)
        cell.text = col
        cell.text_frame.paragraphs[0].runs[0].font.bold = True

    # Data
    for r in range(rows):
        for c in range(cols):
            table.cell(r + 1, c).text = str(df.iloc[r, c])

    return slide














# ============================================================
# ============================================================
# LOAD DATA (ROBUST + NORMALIZED)
# ============================================================
   
# ============================================================
# LOAD DATA
# ============================================================


 # ============================================================
# LOAD DATA (SIMPLE & SAFE)
# ============================================================
@st.cache_data
def load_data(path):
    df = pd.read_excel(path)
    df.columns = df.columns.str.strip()

    df.rename(columns={
        "Session": COL_MONTH,
        "INST NAME": COL_INST,
        "INST CODE": COL_CODE,
        "TLR(30%)": COL_TLR,
        "PPPR(35%)": COL_PPRPI,
        "DO(35%)": COL_DO,
        "Total Score": COL_TOTAL
    }, inplace=True)

    df[COL_MONTH] = df[COL_MONTH].apply(parse_month)
    df[COL_INST] = df[COL_INST].astype(str).str.upper().str.strip()
    df[COL_ZONE] = df[COL_CODE].apply(derive_zone)

    for c in [COL_TOTAL, COL_TLR, COL_PPRPI, COL_DO]:
        if c in df.columns:
            df[c] = pd.to_numeric(df[c], errors="coerce").fillna(0)

    df[COL_NS] = df[COL_TOTAL] == 0
    df[COL_GRADE] = df[COL_TOTAL].apply(grade_from_score)
    df.loc[df[COL_NS], COL_GRADE] = "NS"

    df = df.sort_values([COL_MONTH, COL_TOTAL], ascending=[True, False])
    df[COL_RANK] = df.groupby(COL_MONTH)[COL_TOTAL].rank(
        method="min",
        ascending=False
    )

    return df


# ============================================================
# MAIN APP
# ============================================================
def main():

    components.html(
        """
        <div style="
            padding: 32px 36px;
            margin-bottom: 28px;
            border-radius: 20px;
            background: linear-gradient(90deg, #FFF3E0, #FFE0B2);
            box-shadow: 0 16px 40px rgba(0,0,0,0.12);
            border-left: 10px solid #FF8F00;
            font-family: sans-serif;
        ">

            <h1 style="
                margin: 0;
                font-size: 40px;
                font-weight: 900;
                line-height: 1.15;
                background: linear-gradient(90deg,#FF8F00,#FF6F00);
                -webkit-background-clip: text;
                -webkit-text-fill-color: transparent;
            ">
                🎓 State Institutional Ranking Framework (SIRF) Analytics
            </h1>

            <h2 style="
                margin-top: 12px;
                font-size: 30px;
                font-weight: 700;
                line-height: 1.3;
                color: #6D4C41;
            ">
                Technical Education (Diploma Sector) Uttar Pradesh |
                Institutional Performance Dashboard
            </h2>

        </div>
        """,
        height=220
    )



    df = load_data(DATA_FILE)
    # ============================================================
# LOAD MASTER INSTITUTE LIST (AUTHORITATIVE COUNT)
# ============================================================
    df_master = pd.read_excel(MASTER_FILE)
    
    total_institutes_master = (
        df_master["Institute"]
        .astype(str)
        .str.strip()
        .nunique()
    )
    

    months = sorted(df[COL_MONTH].dropna().unique())
    selected_month = st.sidebar.selectbox(
        "Select Session",
        months,
        format_func=lambda x: x.strftime("%b-%Y")
    )

    df_m = df[df[COL_MONTH] == selected_month]
    zones = sorted(df_m[COL_ZONE].unique())

    tabs = st.tabs([
        "Overview", "Participation", "Zone", "Institute",
        "Parameter-wise", "Red Flags", "Clustering",
        "Map", "XAI", "Export"
    ])




    #for ppt
    
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from io import BytesIO


    
    
    def generate_sirf_ppt(
    session_label,
    kpis,
    fig_grade,
    fig_zone,
    participation_df,
    red_flags_df
):
        prs = Presentation()
    
        # ======================================================
        # SLIDE 1 — TITLE
        # ======================================================
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "State Institutional Ranking Framework (SIRF)"
        slide.placeholders[1].text = (
            "Diploma Sector – Uttar Pradesh\n\n"
            "Evidence-Driven Academic Performance & Inspection Intelligence System\n\n"
            f"Assessment Session: {session_label}"
        )
    
        # ======================================================
        # SLIDE 2 — WHAT THIS SYSTEM PROVIDES
        # ======================================================
        add_text_slide(
            prs,
            "🎯 What This Analytics System Provides",
            "• State-wide visibility of diploma institute performance\n"
            "• Participation tracking across assessment cycles\n"
            "• Zone, district, and institute-level benchmarking\n"
            "• Identification of weak institutes using percentile logic\n"
            "• Inspection prioritisation based on factual indicators\n"
            "• Explainable AI insights explaining loss of marks\n"
            "• Action-oriented outputs for administrators"
        )
    
        # ======================================================
        # SLIDE 3 — EXECUTIVE OVERVIEW
        # ======================================================
        participation_rate = (kpis["participated"] / kpis["total"]) * 100
    
        add_text_slide(
            prs,
            "📊 Executive Snapshot — Current Assessment Cycle",
            f"• Total Institutes under review: {kpis['total']}\n"
            f"• Institutes submitting valid data: {kpis['participated']} "
            f"({participation_rate:.1f}%)\n"
            f"• Institutes skipped / not submitted: {kpis['not_participated']}\n"
            f"• Average score of participating institutes: {kpis['avg_score']:.2f}\n\n"
            "Interpretation:\n"
            "• Participation reflects governance compliance\n"
            "• Average score reflects academic health of the sector"
        )
    
        # ======================================================
        # SLIDE 4 — GRADE DISTRIBUTION
        # ======================================================
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "🎓 Academic Quality Distribution (Grades)"
    
        img = save_fig(fig_grade, "grade_distribution")
        slide.shapes.add_picture(img, Inches(0.8), Inches(1.6), width=Inches(8.5))
    
        # ======================================================
        # SLIDE 5 — ZONE-WISE PERFORMANCE
        # ======================================================
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "🗺️ Zone-wise Average Performance"
    
        img = save_fig(fig_zone, "zone_performance")
        slide.shapes.add_picture(img, Inches(0.8), Inches(1.6), width=Inches(8.5))
    
        # ======================================================
        # SLIDE 6 — PARTICIPATION BEHAVIOUR
        # ======================================================
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "🏫 Institute-wise Participation Behaviour"
    
        rows, cols = participation_df.shape
        table = slide.shapes.add_table(
            rows + 1, cols,
            Inches(0.5), Inches(1.6),
            Inches(9), Inches(4.5)
        ).table
    
        for c, col in enumerate(participation_df.columns):
            table.cell(0, c).text = col
    
        for r in range(rows):
            for c in range(cols):
                table.cell(r + 1, c).text = str(participation_df.iloc[r, c])
    
        # ======================================================
        # SLIDE 7 — RED FLAG INTELLIGENCE
        # ======================================================
        if red_flags_df.empty:
            red_flag_text = (
                "• No institute shows abnormal rank movement (±50)\n\n"
                "Governance Signal:\n"
                "• Data reporting is stable\n"
                "• No immediate physical verification required"
            )
        else:
            red_flag_text = (
                f"• Institutes flagged for verification: {len(red_flags_df)}\n\n"
                "Risk Interpretation:\n"
                "• Sudden rank change may indicate data anomaly\n"
                "• Physical verification recommended"
            )
    
        add_text_slide(
            prs,
            "🚨 Abnormal Rank Movement Detection",
            red_flag_text
        )
    
        # ======================================================
        # SLIDE 8 — EXPLAINABLE AI
        # ======================================================
        add_text_slide(
            prs,
            "🤖 Explainable AI (XAI) — Why Institutes Lose Marks",
            "• Identifies which parameter impacts scores most\n"
            "• Explains why an institute failed to reach full marks\n"
            "• Shows exact marks lost in TLR, PPRPI, and DO\n\n"
            "Strength:\n"
            "• Transparent\n"
            "• Auditable\n"
            "• Policy-defensible"
        )
    
        # ======================================================
        # SLIDE 9 — FROM INSIGHT TO ACTION
        # ======================================================
        add_text_slide(
            prs,
            "🛠️ From Analytics to Administrative Action",
            "Dashboard Insight → Administrative Action\n\n"
            "• Bottom 25% institutes → Academic hand-holding\n"
            "• FSR deficit → Faculty deployment planning\n"
            "• Rank jump ±50 → Physical inspection\n"
            "• Weak DO scores → Placement & industry linkage focus\n"
            "• Critical clusters → Priority monitoring list"
        )
    
        # ======================================================
        # SLIDE 10 — WHY THIS SYSTEM IS DIFFERENT
        # ======================================================
        add_text_slide(
            prs,
            "✅ Why SIRF Analytics Is a Governance Tool",
            "• No subjective judgement\n"
            "• No manual shortlisting\n"
            "• Fully data-driven & rule-based\n"
            "• Scales across sessions, zones, districts\n"
            "• Supports NEP-aligned outcome monitoring\n\n"
            "Objective: Improve institutions, not penalize them"
        )
    
        # ======================================================
        # EXPORT
        # ======================================================
        output = BytesIO()
        prs.save(output)
        output.seek(0)
        return output
    
    
        

    # ========================================================
    # OVERVIEW
    # ========================================================
    # ========================================================
# ========================================================
# TAB 0 — EXECUTIVE OVERVIEW (CLEAN & AUTHORITATIVE)
# ========================================================
    with tabs[0]:
        st.markdown("## 📊 Overview")
    
        # --------------------------------------------------
        # KPI COMPUTATION (SINGLE CATEGORY LOGIC)
        # --------------------------------------------------
    
        # Normalize institute names once (safety)
        df_m[COL_INST] = df_m[COL_INST].astype(str).str.strip().str.upper()
    
        # ✅ Participated = institutes with valid submission
        participated = df_m.loc[df_m[COL_TOTAL] > 0, COL_INST].nunique()
    
        # ✅ Not Participated = NS + missing (collapsed category)
        not_participated = total_institutes_master - participated
    
        # ✅ Average score of participated institutes only
        avg_score = df_m.loc[df_m[COL_TOTAL] > 0, COL_TOTAL].mean()
    
        # Percentages
        participated_pct = round((participated / total_institutes_master) * 100, 1)
        not_participated_pct = round((not_participated / total_institutes_master) * 100, 1)
    
        # --------------------------------------------------
        # PREMIUM KPI CARDS
        # --------------------------------------------------
        c1, c2, c3, c4 = st.columns(4)
    
        c1.metric(
            "🏫 Total Institutes",
            total_institutes_master,
            help="Authoritative count from Master Institute List"
        )
    
        c2.metric(
            "✅ Participated",
            participated,
            delta=f"{participated_pct}%",
            help="Institutes that submitted valid data in this assessment cycle"
        )
    
        c3.metric(
            "❌ Not Participated",
            not_participated,
            delta=f"-{not_participated_pct}%",
            help="Institutes not submitting data or missing from this cycle"
        )
    
        c4.metric(
            "📈 Avg Total Score",
            f"{avg_score:.2f}",
            help="Average score of participating institutes only"
        )
    
        st.markdown("---")
            


    
               # ==================================================
        # 🎓 GRADE DISTRIBUTION (INFORMATIONAL)
        # ==================================================
        st.markdown("### 🎓 Grade Distribution")
        
        # Grades ONLY from participated institutes
        grade_counts = (
            df_m.loc[df_m[COL_TOTAL] > 0, COL_GRADE]
            .value_counts()
        )
        
        # 🔥 REUSE flash-card calculation
        grade_df = pd.DataFrame({
            "Grade": ["A", "B", "C", "D", "NP"],
            "Institutes": [
                grade_counts.get("A", 0),
                grade_counts.get("B", 0),
                grade_counts.get("C", 0),
                grade_counts.get("D", 0),
                not_participated   # ✅ SAME AS KPI CARD
            ]
        })
        
        fig_grade = px.bar(
            grade_df,
            x="Grade",
            y="Institutes",
            color="Grade",
            text="Institutes",
            color_discrete_map=GRADE_COLORS
        )
        
        fig_grade.update_traces(
            textposition="outside",
            cliponaxis=False,      # 🔥 MOST IMPORTANT
            textfont_size=20
        )
        
        fig_grade.update_layout(
            yaxis_title="Number of Institutes",
            xaxis_title="Grade",
            plot_bgcolor="white",
            paper_bgcolor="white",
        
            # 🔥 EXTRA TOP SPACE SO TEXT IS NOT CUT
            margin=dict(t=120, b=50, l=40, r=40),
        
            # 🔥 AUTO SPACE FOR OUTSIDE TEXT
            yaxis=dict(automargin=True)
        )
        
        st.plotly_chart(fig_grade, use_container_width=True)


        # ==================================================
        # 🗺️ ZONE-WISE AVERAGE PERFORMANCE
        # ==================================================
        st.markdown("### 🗺️ Zone-wise Average Performance")
    
        zone_avg = (
            df_m.loc[df_m[COL_TOTAL] > 0]
            .groupby(COL_ZONE, as_index=False)[COL_TOTAL]
            .mean()
            .round(2)
        )
    
        figz = px.bar(
            zone_avg,
            x=COL_ZONE,
            y=COL_TOTAL,
            color=COL_ZONE,
            text=COL_TOTAL,
            color_discrete_map=ZONE_COLORS_DISTINCT
        )
    
        figz.update_traces(
            texttemplate="%{text:.2f}",
            textposition="outside",
            cliponaxis=False,          # 🔥 KEY FIX
            textfont_size=20,
            marker_line_width=1.2,
            marker_line_color="rgba(0,0,0,0.4)"
        )
        
        figz.update_layout(
            showlegend=False,
            yaxis_title="Average Total Score",
            xaxis_title="Zone",
            plot_bgcolor="white",
            paper_bgcolor="white",
        
            # 🔥 EXTRA TOP SPACE
            margin=dict(t=120, b=50, l=40, r=40),
        
            yaxis=dict(
                automargin=True
            )
        )

    
        st.plotly_chart(figz, use_container_width=True)
    
        # ---------------- TABLE ----------------
        zone_avg.index = zone_avg.index + 1
    
        st.dataframe(
            zone_avg.style
                .format({COL_TOTAL: "{:.2f}"})
                .set_properties(**{
                    "text-align": "center",
                    "font-size": "18px"
                })
                .set_table_styles([
                    {"selector": "th", "props": [
                        ("text-align", "center"),
                        ("font-size", "20px")
                    ]}
                ]),
            use_container_width=True
        )


    # ========================================================
    # PARTICIPATION
    # ========================================================
    # =====================================================
# 🏢 INSTITUTE-WISE PARTICIPATION SUMMARY
# =====================================================
# ========================================================
# TAB 1 — INSTITUTE-WISE PARTICIPATION SUMMARY (AUTHORITATIVE)
# ========================================================
# ========================================================
# TAB 1 — INSTITUTE-WISE PARTICIPATION SUMMARY
# ========================================================
    with tabs[1]:
    
        st.markdown("### 🏢 Institute-wise Participation Summary")
        st.caption(
            
            "**Not Submitted + Not Participated are both counted as Skipped.**"
        )
    
        # --------------------------------------------------
        # 1️⃣ CLEAN MASTER (AUTHORITATIVE)
        # --------------------------------------------------
        df_master_clean = df_master.copy()
    
        # Extract 4-digit institute code
        df_master_clean["Inst_Code"] = (
            df_master_clean["Institute"]
            .astype(str)
            .str.extract(r"(\d{4})")
        )
    
        df_master_clean = (
            df_master_clean
            .dropna(subset=["Inst_Code"])
            .assign(
                Inst_Code=lambda x: x["Inst_Code"].astype(str).str.strip(),
                Institute=lambda x: x["Institute"].astype(str).str.upper().str.strip()
            )
            .drop_duplicates(subset=["Inst_Code"])
            .reset_index(drop=True)
        )
    
        # --------------------------------------------------
        # 2️⃣ CLEAN DATA (CRITICAL NORMALIZATION)
        # --------------------------------------------------
        df_data = df.copy()
    
        df_data["Inst_Code"] = df_data["Inst_Code"].astype(str).str.strip()
        df_data[COL_MONTH] = pd.to_datetime(df_data[COL_MONTH]).dt.normalize()
    
        # --------------------------------------------------
        # 3️⃣ GLOBAL SESSIONS (ONLY FROM DATA)
        # --------------------------------------------------
        global_sessions = sorted(df_data[COL_MONTH].dropna().unique())
        total_sessions = len(global_sessions)
    
        # --------------------------------------------------
        # 4️⃣ MASTER × SESSION GRID
        # --------------------------------------------------
        grid = pd.MultiIndex.from_product(
            [df_master_clean["Inst_Code"], global_sessions],
            names=["Inst_Code", COL_MONTH]
        ).to_frame(index=False)
    
        grid[COL_MONTH] = pd.to_datetime(grid[COL_MONTH]).dt.normalize()
    
        # --------------------------------------------------
        # 5️⃣ MERGE ACTUAL DATA
        # --------------------------------------------------
        merged = grid.merge(
            df_data[["Inst_Code", COL_MONTH, COL_TOTAL]],
            on=["Inst_Code", COL_MONTH],
            how="left"
        )
    
        # --------------------------------------------------
        # 6️⃣ PARTICIPATION LOGIC
        # --------------------------------------------------
        merged["Participated"] = merged[COL_TOTAL].fillna(0) > 0
    
        # --------------------------------------------------
        # 7️⃣ SUMMARY PER INSTITUTE
        # --------------------------------------------------
        summary = (
            merged.groupby("Inst_Code", as_index=False)
            .agg(
                Participated_Sessions=("Participated", "sum")
            )
        )
    
        summary["Skipped_Sessions"] = total_sessions - summary["Participated_Sessions"]
    
        summary["Participation (%)"] = (
            summary["Participated_Sessions"] / total_sessions * 100
        ).round(1)
    
        # --------------------------------------------------
        # 8️⃣ STATUS LOGIC
        # --------------------------------------------------
        def status_logic(skipped):
            if skipped >= 3:
                return "🔴 Repeated Non-Participation"
            elif skipped == 2:
                return "🟠 Requires Immediate Follow-up"
            elif skipped == 1:
                return "🟡 Irregular Participation"
            else:
                return "🟢 Regular Participation"
    
        summary["Status"] = summary["Skipped_Sessions"].apply(status_logic)
    
        # --------------------------------------------------
        # 9️⃣ JOIN INSTITUTE NAME
        # --------------------------------------------------
        summary = summary.merge(
            df_master_clean[["Inst_Code", "Institute"]],
            on="Inst_Code",
            how="left"
        )
    
        # --------------------------------------------------
        # 🔟 SORT + SERIAL NUMBER
        # --------------------------------------------------
        summary = summary.sort_values(
            ["Skipped_Sessions", "Participation (%)", "Institute"],
            ascending=[False, True, True]
        ).reset_index(drop=True)
    
        summary.insert(0, "S.No", summary.index + 1)
    
        # --------------------------------------------------
        # 1️⃣1️⃣ DISPLAY TABLE
        # --------------------------------------------------
        st.dataframe(
            summary[
                [
                    "S.No",
                    "Inst_Code",
                    "Institute",
                    "Skipped_Sessions",
                    "Participation (%)",
                    "Status"
                ]
            ],
            use_container_width=True,
            height=560,
            hide_index=True
        )
    
    
    
    
    
    
    
    












  

# ========================================================
    with tabs[2]:
        st.markdown("## 🗺️ Zone Wise Performance Overview")
        st.caption(
            "This section shows how institutes in a selected zone are performing. "
            "You can quickly see participation, average performance, strong institutes, "
            "and institutes that need support."
        )
    
        # Select zone
        zone = st.selectbox("Select Zone to View", zones)
        dz = df_m[df_m[COL_ZONE] == zone].copy()
    
        if dz.empty:
            st.warning("No data available for this zone.")
            st.stop()
    
        # ---------------------------
        # KEY SUMMARY (IN SIMPLE WORDS)
        # ---------------------------
        total_inst = dz[COL_INST].nunique()
        participated = dz[~dz[COL_NS]][COL_INST].nunique()
        avg_score = dz[COL_TOTAL].mean()
        common_grade = dz[COL_GRADE].value_counts().idxmax()
    
        c1, c2, c3, c4 = st.columns(4)
        c1.metric("Total Institutes", total_inst)
        c2.metric("Institutes Participated", participated)
        c3.metric("Average Total Score", f"{avg_score:.2f}")
        c4.metric("Most Institutes Fall in Grade", common_grade)
    
        st.markdown("---")

        # ---------------------------
        # PERFORMANCE MAP (EASY EXPLANATION)
        # ---------------------------
        st.markdown("### 📍 How Institutes Are Performing")
        st.caption(
            "Each circle represents one institute.\n\n"
            "- ➡️ Right side → Better teaching resources (TLR)\n"
            "- ⬆️ Upper side → Better practical & project learning (PPRPI)\n"
            "- 🔵 Bigger circle → Better diploma outcome (DO)\n\n"
            "Institutes in the **top-right with larger circles** are performing the best."
        )
    
        fig = px.scatter(
            dz,
            x=COL_TLR,
            y=COL_PPRPI,
            size=COL_DO,                     # ✅ CORRECTED
            color=COL_GRADE,
            hover_name=COL_INST,
            hover_data={
                COL_TLR: ':.1f',
                COL_PPRPI: ':.1f',
                COL_DO: ':.1f',
                COL_TOTAL: ':.1f'
            },
            size_max=60,
            color_discrete_map=GRADE_COLORS,
            opacity=0.85
        )
    
        # Reference lines (median-based guidance)
        fig.add_vline(
            x=dz[COL_TLR].median(),
            line_dash="dash",
            line_color="gray",
            annotation_text="Average TLR",
            annotation_position="top"
        )
        fig.add_hline(
            y=dz[COL_PPRPI].median(),
            line_dash="dash",
            line_color="gray",
            annotation_text="Average PPRPI",
            annotation_position="right"
        )
    
        fig.update_layout(
            xaxis_title="Teaching, Learning & Resources (TLR)",
            yaxis_title="Practical, Project & Research Performance (PPRPI)",
            plot_bgcolor="white",
            paper_bgcolor="white"
        )
    
        st.plotly_chart(fig, use_container_width=True)
    
        # ---------------------------
        st.markdown("### 🎓 Grade Distribution in This Zone")
        st.caption(
            "This chart shows how many institutes fall under each grade "
            "(A = very good, D = needs serious improvement)."
        )
        
        grade_zone = (
            dz[COL_GRADE]
            .value_counts()
            .reindex(["A", "B", "C", "D", "NS"], fill_value=0)
            .reset_index(name="Number of Institutes")
            .rename(columns={"index": "Grade"})
        )
        
        figg = px.bar(
            grade_zone,
            x="Grade",
            y="Number of Institutes",
            color="Grade",
            text="Number of Institutes",
            color_discrete_map=GRADE_COLORS
        )
        
        y_max = grade_zone["Number of Institutes"].max()
        
        figg.update_traces(
            textposition="outside",
            cliponaxis=False
        )
        
        figg.update_yaxes(
            range=[0, y_max * 1.25]
        )
        
        figg.update_layout(
            margin=dict(t=90),
            uniformtext_minsize=14,
            uniformtext_mode="show"
        )
        
        st.plotly_chart(figg, use_container_width=True)

        # ---------------------------
        # BEST & WORST INSTITUTES
        # ---------------------------
        st.markdown("### 🏆 Top 5 Performers & 🚨 Bottom 5 Institutes Needing Support")
        st.caption(
            "Top institutes can be used as role models. "
            "Bottom institutes may need academic, faculty, or placement support."
        )
    
                # -------------------------------------------------
        # CLEAN EXCEL SERIAL NUMBER (RUN ONCE AFTER LOAD)
        # -------------------------------------------------
        dz = dz.drop(
            columns=["S.No", "S No", "Sr No", "Sr. No", "Index", "Unnamed: 0"],
            errors="ignore"
        )
        
        # -------------------------------------------------
        # TOP & BOTTOM 5 (ADMIN VIEW)
        # -------------------------------------------------
        top5 = dz.sort_values(COL_TOTAL, ascending=False).head(5).copy()
        bottom5 = dz.sort_values(COL_TOTAL).head(5).copy()
        
        # Generate fresh serial numbers (1–5)
        top5["S. No."] = range(1, len(top5) + 1)
        bottom5["S. No."] = range(1, len(bottom5) + 1)
        
        # Column order (Serial first)
        display_cols = [
            "S. No.",
            COL_INST,
            COL_TOTAL,
            COL_DO,
            COL_GRADE,
            COL_RANK
        ]
        
        # -------------------------------------------------
        # DISPLAY (STREAMLIT)
        # -------------------------------------------------
        c1, c2 = st.columns(2)
        
        with c1:
            st.markdown("#### 🏆 Top 5 Institutes")
            st.dataframe(
                top5[display_cols],
                use_container_width=True,
                hide_index=True
            )
        
        with c2:
            st.markdown("#### 🚨 Bottom 5 Institutes (Need Attention)")
            st.dataframe(
                bottom5[display_cols],
                use_container_width=True,
                hide_index=True
            )

    
        # ---------------------------
        # COMPLETE LIST
        # ---------------------------
        all_df = dz.sort_values(COL_RANK).copy()
        all_df["S. No."] = range(1, len(all_df) + 1)

        st.markdown("### 📋 All Institutes in This Zone")
        st.caption(
            "Complete list of institutes in this zone, sorted by rank "
            "(Rank 1 = best performer)."
        )
        
        st.dataframe(
            all_df[
                [
                    "S. No.",
                    COL_INST,
                    COL_TOTAL,
                    COL_TLR,
                    COL_PPRPI,
                    COL_DO,
                    COL_GRADE,
                    COL_RANK
                ]
            ],
            use_container_width=True,
            height=520,
            hide_index=True
        )

    
        



 


  # ========================================================
# INSTITUTE (DETAILED, HUMAN-UNDERSTANDABLE VIEW)
# ========================================================
   
        # ========================================================
# TAB 3 — INSTITUTE (CLEAR, EXPLAINABLE, ACTION-ORIENTED)
# ========================================================
    with tabs[3]:
        st.markdown("## 🏢 Institute Performance Dashboard")
        st.caption(
            "This section shows the complete performance story of a single institute — "
            "current standing, past trends, strengths, and areas needing attention."
        )
    
        # ---------------------------
        # SELECT INSTITUTE
        # ---------------------------
        inst = st.selectbox(
            "Select Institute",
            sorted(df[COL_INST].unique())
        )
    
        hist = df[df[COL_INST] == inst].sort_values(COL_MONTH)
    
        if hist.empty:
            st.warning("No data available for this institute.")
            st.stop()
    
        latest = hist.iloc[-1]

        # ---------------------------
        # CURRENT STATUS (AT A GLANCE)
        # ---------------------------
        st.markdown("### 📌 Institute Performance Overview)")
    
        c1, c2, c3, c4 = st.columns(4)
        c1.metric("Assessment Month (Latest) ", latest[COL_MONTH].strftime("%b-%Y"))
        c2.metric("Overall Score", f"{latest[COL_TOTAL]:.2f}")
        c3.metric("Grade", latest[COL_GRADE])
        c4.metric("Rank", int(latest[COL_RANK]))
    
        st.markdown("---")
    
        # ---------------------------
        # PERFORMANCE TREND (VISUAL + SIGNALS)
        # ---------------------------
        st.markdown("### 📈 Performance Trend Over Time")
        st.caption(
            "Each line shows how the institute has progressed across sessions.\n\n"
            "- **Overall Performance** → Total Score\n"
            "- **Teaching & Resources (TLR)**\n"
            "- **Practical, Project & Research Impact (PPRPI)**\n"
            "- **Diploma Outcomes (DO)**\n\n"
            "Visual signals:\n"
            "🟢 Green triangle → Overall improvement\n"
            "🔴 Red dot → Diploma outcome decline"
        )

        # Prepare plotting data
        plot_df = hist.copy()
        plot_df["Session"] = plot_df[COL_MONTH].dt.strftime("%b-%y")
    
        fig = px.line(
            plot_df,
            x="Session",
            y=[COL_TOTAL, COL_TLR, COL_PPRPI, COL_DO],
            markers=True,
            labels={
                "value": "Score",
                "variable": "Performance Area",
                "Session": "Session"
            }
        )
    
        # Human-friendly legend names
        legend_map = {
            COL_TOTAL: "Overall Performance (Total Score)",
            COL_TLR: "Teaching, Learning & Resources (TLR)",
            COL_PPRPI: "Practical, Project & Research Impact (PPRPI)",
            COL_DO: "Diploma Outcomes (DO)"
        }
    
        fig.for_each_trace(
            lambda t: t.update(name=legend_map.get(t.name, t.name))
        )

        # 🔴 Red dot where DO drops
        do_drop = plot_df[plot_df[COL_DO].diff() < 0]
        if not do_drop.empty:
            fig.add_scatter(
                x=do_drop["Session"],
                y=do_drop[COL_DO],
                mode="markers",
                marker=dict(color="red", size=14),
                name="DO Decline"
            )
    
        # 🟢 Green arrow where Total improves
        improve = plot_df[plot_df[COL_TOTAL].diff() > 0]
        if not improve.empty:
            fig.add_scatter(
                x=improve["Session"],
                y=improve[COL_TOTAL],
                mode="markers",
                marker=dict(color="green", size=16, symbol="triangle-up"),
                name="Overall Improvement"
            )
    
        fig.update_layout(
            xaxis_title="Session",
            yaxis_title="Score",
            legend_title_text="Performance Indicators",
            plot_bgcolor="white",
            paper_bgcolor="white",
            hovermode="x unified",
            xaxis=dict(tickangle=-30)
        )
    
        st.plotly_chart(fig, use_container_width=True)

        # ---------------------------
        # AUTO INTERPRETATION (NO THINKING REQUIRED)
        # ---------------------------
        st.markdown("### 🧠 Performance Interpretation")
    
        if len(plot_df) >= 2:
            total_change = plot_df[COL_TOTAL].iloc[-1] - plot_df[COL_TOTAL].iloc[-2]
            do_change = plot_df[COL_DO].iloc[-1] - plot_df[COL_DO].iloc[-2]
    
            if total_change > 0 and do_change >= 0:
                st.success(
                    "🟢 **Performance Improving** — Overall score and diploma outcomes are both improving."
                )
            elif total_change > 0 and do_change < 0:
                st.warning(
                    "⚠️ **Mixed Signal** — Overall score improved, but diploma outcomes declined. "
                    "Teaching quality may not be translating into student results."
                )
            elif total_change < 0:
                st.error(
                    "🔴 **Performance Declining** — Overall performance has dropped compared to last session."
                )
            else:
                st.info(
                    "ℹ️ **Performance Stable** — No major change from the previous session."
                )
        else:
            st.info("Not enough historical data to determine a trend.")
    
        st.markdown("---")
    
        # ---------------------------
        # SESSION-WISE DETAILS (FOR REVIEW & AUDIT)
        # ---------------------------
        
        display_df = hist.copy()
        
        # 1️⃣ Force clean Month (drop invalid / empty)
        display_df[COL_MONTH] = pd.to_datetime(
            display_df[COL_MONTH], errors="coerce"
        )
        
        # 2️⃣ Keep ONLY valid rows (STRICT FILTER)
        display_df = display_df[
            display_df[COL_MONTH].notna() &
            (
                display_df[COL_TOTAL].notna() |
                display_df[COL_TLR].notna() |
                display_df[COL_PPRPI].notna() |
                display_df[COL_DO].notna()
            )
        ]
        
        # 3️⃣ Sort chronologically
        display_df = display_df.sort_values(COL_MONTH).reset_index(drop=True)
        
        # 4️⃣ Format month for display
        display_df[COL_MONTH] = display_df[COL_MONTH].dt.strftime("%b-%Y")
        
        # 5️⃣ Generate correct serial number
        display_df.insert(0, "S. No.", range(1, len(display_df) + 1))
        
        # ---------------------------
        # DISPLAY
        # ---------------------------
        st.markdown("### 📋 Assessment Month-wise Detailed Performance")
        st.caption(
            "This table shows all Assessment Monthn-wise scores and ranks. "
            "Useful for audits, reviews, and official reporting."
        )
        
        st.dataframe(
            display_df[
                [
                    "S. No.",
                    COL_MONTH,
                    COL_TOTAL,
                    COL_TLR,
                    COL_PPRPI,
                    COL_DO,
                    COL_GRADE,
                    COL_RANK
                ]
            ],
            use_container_width=True,
            height=520,
            hide_index=True
        )





        
       
#-------------------------------------------------------------------------------------
  
    # ========================================================
# TAB 4 — PARAMETER-WISE DIAGNOSTIC, ROOT CAUSE & ACTION
# # ========================================================
# TAB 4 — PARAMETER-WISE DIAGNOSTIC & ROOT CAUSE (ADMIN READY)
# ========================================================
    with tabs[4]:
    
        # ====================================================
        # TITLE & CONTEXT (ADMIN FRIENDLY)
        # ====================================================
        st.markdown("## 🛠️ Parameter Wise Analysis")
        st.caption(
            "This section identifies **exact operational weaknesses** "
            "responsible for low scores. The analysis is:\n\n"
            "• Based only on **submitted data**\n"
            "• Uses **policy thresholds**, not assumptions\n"
            "• Suitable for **inspection, review & corrective action**"
        )
    
        # --------------------------------------------------
        # PARAMETER SELECTION
        # --------------------------------------------------
        PARAM_MAP = {
            COL_TLR: "Teaching, Learning & Resources (TLR)",
            COL_PPRPI: "Practical, Project & Research Impact (PPRPI)",
            COL_DO: "Diploma Outcomes (DO)"
        }
    
        param = st.selectbox(
            "Select Performance Dimension",
            list(PARAM_MAP.keys()),
            format_func=lambda x: PARAM_MAP[x]
        )
        param_name = PARAM_MAP[param]
    
        # --------------------------------------------------
        # GAP / DEFICIENCY VALUES (CATEGORICAL)
        # --------------------------------------------------
        GAP_VALUES = {
            "NO", "0", "0%",
            "LESS THAN 50%", "LESS THAN 70%",
            "NA", "N/A", "NONE",
            "NOT AVAILABLE", "", "NULL"
        }
    
        # --------------------------------------------------
        # CONSIDER ONLY PARTICIPATED INSTITUTES
        # --------------------------------------------------
        valid = df_m[df_m[COL_TOTAL] > 0].copy()
        if valid.empty:
            st.warning("⚠️ No institutes participated in ranking for this session.")
            st.stop()
    
        # --------------------------------------------------
        # BOTTOM 25% — FACTUAL IDENTIFICATION
        # --------------------------------------------------
        q25 = valid[param].quantile(0.25)
        focus = valid[valid[param] <= q25].copy()
    
        st.markdown(f"### 🔍 Bottom 25% Institutes — **{param_name}**")
        st.write(f"**Score Threshold:** {q25:.2f}")
    
        focus_tbl = (
            focus[[COL_INST, COL_ZONE, "District", param]]
            .sort_values(param)
            .reset_index(drop=True)
        )
    
        st.dataframe(
            df_index(focus_tbl),
            use_container_width=True,
            height=360,
            hide_index=True
        )
    
        # ====================================================
        # ROOT CAUSE — SUB PARAMETER DEFICIENCY
        # ====================================================
        st.markdown("## 🧠 Root Cause Analysis (Parameter-wise Deficiency Count)")
    
        # ---------------- TLR ----------------
        if param == COL_TLR:
            sub_params = [
                "faculty student ration",
                "b.tech/B.E./M.A./M.Sc",
                "m.tech/M.E./NET/SLET/SET",
                "P.HD",
                "Experience < 2 Years",
                "2 < = Experience <= 5 Years",
                "5 < Experience <= 10 Years",
                "Experience > 10 Years",
                "Monthly Average Students Attendance",
    
                # Infrastructure
                "RO Drinking Water",
                "electricity upto 16 Hrs/Day",
                "Backup Availability",
                "Toilet in Administrative Block",
                "Toilets for Boys",
                "Toilets for Girls",
                "Sweeper",
                "Internet",
                "Wi-Fi",
                "Smart Class",
    
                # Labs & Library
                "Number of Labs/Workshops required",
                "Total no. of Labs/Workshops available in the institute",
                "Total no. of Computers required",
                "Total no. of Computers available",
                "Total no. of Computers working",
                "E-Library with E-Granthalaya"
            ]
    
        # ---------------- PPRPI ----------------
        elif param == COL_PPRPI:
            sub_params = [
                "MOOCs Courses,Swayam Portal, NPTEL, NITTT Modules",
                "Departmental Training",
                "Other than departmental training",
                "Faculty pursuing Higher Education",
    
                "Project to faculty engagement ratio in each discipline",
                "Total number of Jobs/Practicals Scheduled",
                "Group size formed in the institute for the project work",
                "Monthly Progress Report of project work/models",
    
                "Institute-Industry Relationship Index",
                "Project Exhibition Index",
                "Student Appreciation Index",
                "Institute Competitive Index",
    
                "Publications Index",
                "Intellectual Property Right",
                "Workshop Index",
                "Seminar Index"
            ]
    
        # ---------------- DO ----------------
        else:
            sub_params = [
                "On Time Diploma Completion Rate",
                "Percentage of students passed with honours",
                "Placement Success Rate",
                "Percentage of students participated in skill development activities",
                "Companies visited for placement",
                "No of Students placed through campus placement",
                "No of Students placed other than campus placement",
                "No of Students opted higher education",
                "Mock interviews for students conducted"
            ]
    
        # --------------------------------------------------
        # DEFICIENCY COMPUTATION
        # --------------------------------------------------
        gap_summary = []
        institute_gaps = []
    
        for col in sub_params:
            if col not in focus.columns:
                continue
    
            # Faculty–Student Ratio (Policy logic)
            if col == "faculty student ration":
                fsr_str = focus[col].fillna("").astype(str)
                ratio = fsr_str.str.extract(r'1\s*:\s*(\d+)', expand=False)
                fsr_val = pd.to_numeric(ratio, errors="coerce")
                mask = fsr_val > 25
    
            else:
                mask = (
                    focus[col]
                    .fillna("")
                    .astype(str)
                    .str.strip()
                    .str.upper()
                    .isin(GAP_VALUES)
                )
    
            affected = mask.sum()
    
            if affected > 0:
                label = SUB_PARAM_LABELS.get(col, col)
    
                gap_summary.append({
                    "Deficient Parameter": label,
                    "Institutes Affected": int(affected)
                })
    
                for _, r in focus[mask].iterrows():
                    institute_gaps.append({
                        "Institute": r[COL_INST],
                        "District": r["District"],
                        "Zone": r[COL_ZONE],
                        "Deficiency": label
                    })
    
        # ====================================================
        # OUTPUTS
        # ====================================================
        if gap_summary:
    
            gap_df = (
                pd.DataFrame(gap_summary)
                .sort_values("Institutes Affected", ascending=False)
                .reset_index(drop=True)
            )
            gap_df.insert(0, "S. No.", range(1, len(gap_df) + 1))
    
            st.markdown("### 🚨 High-Impact Weak Areas")
            st.data_editor(
                gap_df,
                use_container_width=True,
                hide_index=True,
                disabled=True
            )
    
            st.markdown("### 📋 Institute-wise Evidence (Inspection Ready)")
            inst_gap_df = pd.DataFrame(institute_gaps).drop_duplicates()
            inst_gap_df.insert(0, "S. No.", range(1, len(inst_gap_df) + 1))
    
            st.data_editor(
                inst_gap_df,
                use_container_width=True,
                hide_index=True,
                disabled=True
            )
    
            st.download_button(
                "📤 Download Deficiency Evidence (CSV)",
                inst_gap_df.to_csv(index=False).encode(),
                file_name="institute_wise_deficiency_evidence.csv",
                mime="text/csv"
            )
    
        else:
            st.success("✅ No measurable deficiencies detected for this parameter.")


    # ========================================================
    # RED FLAGS (RANK JUMP ≥20)
    # ========================================================
    # ========================================================
# TAB 5 — RED FLAGS (±50 Rank Movement)
# ========================================================
    with tabs[5]:
        st.markdown("## 🚨 Red Flag Analysis — Rank Movement (±50)")
        st.caption(
            "This section flags institutes with **abnormal rank movement** "
            "between the current and previous **assessment month**.\n\n"
            "⚠️ Movements of **±50 ranks or more** require **physical verification** "
            "to validate data accuracy, reporting changes, or operational issues."
        )
    
        # --------------------------------------------------
        # PREVIOUS vs CURRENT RANK
        # --------------------------------------------------
        prev_rank = (
            df[df[COL_MONTH] < selected_month]
            .groupby(COL_INST)[COL_RANK]
            .last()
        )
    
        curr_rank = df_m.set_index(COL_INST)[COL_RANK]
    
        # Align institutes
        rank_compare = (
            pd.concat([prev_rank, curr_rank], axis=1, keys=["Previous Rank", "Current Rank"])
            .dropna()
            .reset_index()
        )
    
        # Rank movement (positive = improvement, negative = decline)
        rank_compare["Rank Change"] = (
            rank_compare["Previous Rank"] - rank_compare["Current Rank"]
        )
    
        # --------------------------------------------------
        # APPLY ±50 THRESHOLD
        # --------------------------------------------------
        red_flags = rank_compare[
            rank_compare["Rank Change"].abs() >= 50
        ].copy()
    
        if red_flags.empty:
            st.success("✅ No institute shows abnormal rank movement (±50) in this assessment month.")
            st.stop()
    
        # --------------------------------------------------
        # CLASSIFY MOVEMENT TYPE (FACTUAL)
        # --------------------------------------------------
        def movement_type(x):
            if x >= 50:
                return "🟢 Sharp Improvement (≥ +50)"
            elif x <= -50:
                return "🔴 Sharp Decline (≤ −50)"
            return "—"
    
        red_flags["Movement Type"] = red_flags["Rank Change"].apply(movement_type)
        red_flags["Action Required"] = "🔍 Physical Verification Required"
    
        # --------------------------------------------------
        # FINAL DISPLAY TABLE
        # --------------------------------------------------
        red_flags = red_flags[
            [
                COL_INST,
                "Previous Rank",
                "Current Rank",
                "Rank Change",
                "Movement Type",
                "Action Required"
            ]
        ]
    
        red_flags = df_index(red_flags)   # ✅ S.No from 1
        
                # --------------------------------------------------
        # SUMMARY COUNTS — ABRUPT UP & DOWN MOVEMENTS
        # --------------------------------------------------
        up_count = (red_flags["Rank Change"] >= 50).sum()
        down_count = (red_flags["Rank Change"] <= -50).sum()
        
        c1, c2 = st.columns(2)
        
        c1.metric(
            label="⬆️ Abrupt Rank Improvement (≥ +50)",
            value=up_count
        )
        
        c2.metric(
            label="⬇️ Abrupt Rank Decline (≤ −50)",
            value=down_count
        )

        st.dataframe(
            red_flags,
            use_container_width=True,
            height=420
        )
    
        # --------------------------------------------------
        # ADMIN INTERPRETATION (FACT-ONLY)
        # --------------------------------------------------
        st.markdown("### 🧭 Interpretation ")
        st.markdown(
            """
    - 🔴 **Sharp Decline (≤ −50 ranks)**  
      May indicate missing data, reporting errors, or operational disruption  
      → **Mandatory physical inspection recommended**
    
    - 🟢 **Sharp Improvement (≥ +50 ranks)**  
      May indicate data correction, delayed reporting, or sudden performance change  
      → **Verification required to confirm sustainability**
            """
        )
    
        # --------------------------------------------------
        # EXPORT FOR FIELD VERIFICATION
        # --------------------------------------------------
        st.download_button(
            "📤 Download Red-Flag Institute List (For Physical Verification)",
            red_flags.to_csv(index=False).encode(),
            file_name="red_flag_rank_movement_physical_verification.csv",
            mime="text/csv"
        )


    # ========================================================
    # CLUSTERING
    # ========================================================
    # ========================================================
# CLUSTERING — PERFORMANCE SEGMENTATION (ADMIN FRIENDLY)
# ========================================================
    with tabs[6]:
        st.markdown("## 🧩 Institute Performance Clustering")
        st.caption(
            "Institutes are grouped based on **overall performance pattern** "
            "(Total Score, TLR, PPRPI, DO).\n\n"
            "Each cluster represents institutes with **similar academic and outcome behaviour**."
        )

        # --------------------------------------------------
        # PREPARE DATA
        # --------------------------------------------------
        cluster_df = df_m[df_m[COL_TOTAL] > 0].copy()
    
        X = cluster_df[[COL_TOTAL, COL_TLR, COL_PPRPI, COL_DO]]
        X_scaled = StandardScaler().fit_transform(X)
    
        # --------------------------------------------------
        # K-MEANS CLUSTERING
        # --------------------------------------------------
        km = KMeans(n_clusters=4, random_state=42, n_init=10)
        cluster_df["Cluster"] = km.fit_predict(X_scaled)
    
        # --------------------------------------------------
        # PCA FOR VISUALIZATION
        # --------------------------------------------------
        pcs = PCA(n_components=2, random_state=42).fit_transform(X_scaled)
        cluster_df["PC1"] = pcs[:, 0]
        cluster_df["PC2"] = pcs[:, 1]
    
        # --------------------------------------------------
        # CLUSTER PROFILING (MEAN SCORES)
        # --------------------------------------------------
        cluster_profile = (
            cluster_df
            .groupby("Cluster")[[COL_TOTAL, COL_TLR, COL_PPRPI, COL_DO]]
            .mean()
            .round(2)
            .reset_index()
        )
    
        # Rank clusters by Total Score
        cluster_profile = cluster_profile.sort_values(COL_TOTAL, ascending=False)
        cluster_profile["Performance Band"] = [
            "🟢 High Performing",
            "🟡 Above Average",
            "🟠 Needs Improvement",
            "🔴 Critical"
        ]
    
        # Map labels back to main df
        cluster_label_map = dict(
            zip(cluster_profile["Cluster"], cluster_profile["Performance Band"])
        )
        cluster_df["Performance Band"] = cluster_df["Cluster"].map(cluster_label_map)
    
        # --------------------------------------------------
        # KPI SUMMARY (ADMIN VIEW)
        # --------------------------------------------------
        st.markdown("### 📌 Cluster Summary (At a Glance)")
    
        c1, c2, c3, c4 = st.columns(4)
        c1.metric("🟢 High Performing", (cluster_df["Performance Band"] == "🟢 High Performing").sum())
        c2.metric("🟡 Above Average", (cluster_df["Performance Band"] == "🟡 Above Average").sum())
        c3.metric("🟠 Needs Improvement", (cluster_df["Performance Band"] == "🟠 Needs Improvement").sum())
        c4.metric("🔴 Critical", (cluster_df["Performance Band"] == "🔴 Critical").sum())
    
        st.markdown("---")
    
        # --------------------------------------------------
        # VISUAL CLUSTER MAP
        # --------------------------------------------------
        st.markdown("### 📍 Performance Cluster Map")
        st.caption(
            "Each dot is an institute.\n\n"
            "- Dots close together → similar performance pattern\n"
            "- Right/top direction → stronger performance\n"
            "- 🔴 Critical cluster → highest priority for hand-holding"
        )
    
        fig = px.scatter(
            cluster_df,
            x="PC1",
            y="PC2",
            color="Performance Band",
            size=COL_TOTAL,
            hover_name=COL_INST,
            hover_data={
                COL_TOTAL: ':.1f',
                COL_TLR: ':.1f',
                COL_PPRPI: ':.1f',
                COL_DO: ':.1f',
                COL_ZONE: True
            },
            color_discrete_map={
                "🟢 High Performing": "#2E7D32",
                "🟡 Above Average": "#F9A825",
                "🟠 Needs Improvement": "#FB8C00",
                "🔴 Critical": "#C62828"
            },
            opacity=0.85
        )
    
        fig.update_layout(
            xaxis_title="Performance Pattern Axis 1",
            yaxis_title="Performance Pattern Axis 2",
            plot_bgcolor="white",
            paper_bgcolor="white",
            legend_title_text="Cluster Category"
        )
    
        st.plotly_chart(fig, use_container_width=True)
    
        # --------------------------------------------------
        # CLUSTER-WISE INSTITUTE LIST (ACTIONABLE)
        # --------------------------------------------------
        st.markdown("### 📋 Institutes by Cluster (For Action)")
    
        cluster_table = (
            cluster_df[
                [
                    COL_INST, "District", COL_ZONE,
                    COL_TOTAL, COL_TLR, COL_PPRPI, COL_DO,
                    "Performance Band"
                ]
            ]
            .sort_values(["Performance Band", COL_TOTAL], ascending=[True, False])
            .reset_index(drop=True)
        )
        cluster_table.insert(0, "S.No", range(1, len(cluster_table) + 1))
    
        st.dataframe(cluster_table, use_container_width=True, height=520,
            hide_index=True,)
    
      
    

    # ========================================================
    # MAP
    # ========================================================
# ========================================================
# TAB 7 — DISTRICT INSPECTION PRIORITY MAP (UP GEOJSON)
   # ========================================================
    # TAB 7 — DISTRICT INSPECTION PRIORITY (MAP + SCORE LOGIC)
    # ========================================================
    with tabs[7]:
        #st.markdown("## 🗺️ District Inspection Priority — Uttar Pradesh")
    
    
        # ====================================================
        # 📊 DATA & GEOJSON LOADING
        # ====================================================
        import json, folium
        from streamlit_folium import st_folium
    
        GEOJSON_PATH = r"data/uttar_pradesh_districts.geojson"
    
        with open(GEOJSON_PATH, "r", encoding="utf-8") as f:
            up_geojson = json.load(f)
    
        # ====================================================
        # 🔧 NORMALIZATION + ALIAS (CRITICAL)
        # ====================================================
        def normalize_dist(x):
            if pd.isna(x):
                return None
        
            x = str(x).upper()
        
            # basic normalization
            x = x.replace("&", "AND")
            x = x.replace("-", " ")
        
            # remove administrative noise (CRITICAL)
            for w in [" DISTRICT", " NAGAR", " CITY"]:
                x = x.replace(w, "")
        
            # collapse multiple spaces
            x = " ".join(x.split())
        
            return x.strip()


        DISTRICT_ALIAS = {

            # ================= Official renamed districts =================
            "ALLAHABAD": "PRAYAGRAJ",
            "FAIZABAD": "AYODHYA",
            "JYOTIBA PHULE": "AMROHA",
            "KANSHIRAM": "KASGANJ",
            "MAHAMAYA": "HATHRAS",
            "SANT RAVIDAS": "BHADOHI",
        
            # ================= Compound / merged districts =================
            "KHERI": "LAKHIMPUR KHERI",
            "LAKHIMPUR": "LAKHIMPUR KHERI",
        
            # ================= Spelling / spacing variants =================
            "BARA BANKI": "BARABANKI",
            "RAE BARELI": "RAEBARELI",
            "BULANDSHAHR": "BULANDSHAHAR",
            "KAUSHAMBI": "KAUSHAMBHI",
            "SIDDHARTHNAGAR": "SIDHARTHNAGAR",
            "SHRAWASTI": "SHRAVASTI",
            "BAGHPAT": "BAGPAT",
            "JAUN PUR": "JAUNPUR",
            "SULTAN PUR": "SULTANPUR",
            "GHAZI PUR": "GHAZIPUR",
        
            # ================= Misspellings / phonetic variants =================
            "BEHRAICH": "BAHRAICH",
            "SANTRAVIDAS": "BHADOHI",
            "KANSHIRAN": "KASGANJ",
            "GAUTAMBUDH": "GAUTAMBUDDHANAGAR",
            "GAUTAM BUDH": "GAUTAMBUDDHANAGAR"
        }

    
        # ====================================================
        # 🗺️ APPLY NORMALIZATION + ALIAS TO GEOJSON
        # ====================================================
        all_districts = []
    
        for feat in up_geojson["features"]:
            raw = normalize_dist(feat["properties"].get("district_name"))
            mapped = DISTRICT_ALIAS.get(raw, raw)
    
            feat["properties"]["district_norm"] = mapped
    
            if mapped:
                all_districts.append(mapped)
    
        # ====================================================
        # 📊 BASE DATA (ONLY PARTICIPATED INSTITUTES)
        # ====================================================
        work = df_m[df_m[COL_TOTAL] > 0].copy()
    
        work["District_norm"] = (
            work["District"]
            .apply(normalize_dist)
            .replace(DISTRICT_ALIAS)
        )
    
        q25 = work[COL_TOTAL].quantile(0.25)
        q50 = work[COL_TOTAL].quantile(0.50)
        q75 = work[COL_TOTAL].quantile(0.75)
    
        FSR_COL = "faculty student ration"
        FSR_DEFICIT = {
            "1:26 TO 1:35",
            "1:36 TO 1:45",
            "1:45 TO 1:50",
            "GREATER THAN 1:50"
        }
    
        prev = (
            df[df[COL_MONTH] < selected_month]
            .groupby(COL_INST)[COL_RANK]
            .last()
        )
        curr = df_m.set_index(COL_INST)[COL_RANK]
        rank_jump = (prev - curr).abs()
        rank_jump = rank_jump[rank_jump >= 50]
    
        # ====================================================
        # 🧮 DISTRICT INSPECTION SCORE
        # ====================================================
        rows = []
    
        for dist in all_districts:
            g = work[work["District_norm"] == dist]
    
            b25 = (g[COL_TOTAL] <= q25).sum()
            b50 = ((g[COL_TOTAL] > q25) & (g[COL_TOTAL] <= q50)).sum()
            b75 = ((g[COL_TOTAL] > q50) & (g[COL_TOTAL] <= q75)).sum()
    
            fsr = (
                g[FSR_COL].astype(str).str.upper().isin(FSR_DEFICIT).sum()
                if not g.empty else 0
            )
    
            rjump = (
                rank_jump.reindex(g[COL_INST]).dropna().shape[0]
                if not g.empty else 0
            )
    
            score = (b25 * 4) + (b50 * 3) + (b75 * 2) + (fsr * 2) + (rjump * 3)
    
            rows.append({
                "District": dist,
                "Bottom 25%": b25,
                "25–50%": b50,
                "50–75%": b75,
                "FSR Deficit": fsr,
                "Rank Jump ±50": rjump,
                "Inspection Score": score
            })
    
        district_df = pd.DataFrame(rows)
    
        # ====================================================
        # 🗺️ CHOROPLETH MAP
        # ====================================================
        st.markdown("## 🗺️ Inspection Priority Map")
    
        m = folium.Map(
            location=[26.85, 80.95],
            zoom_start=6,
            tiles="CartoDB positron"
        )
    
        folium.Choropleth(
            geo_data=up_geojson,
            data=district_df,
            columns=["District", "Inspection Score"],
            key_on="feature.properties.district_norm",
            fill_color="YlOrRd",
            fill_opacity=0.9,
            line_opacity=0.4,
            nan_fill_color="#E0E0E0",
            legend_name="Inspection Priority Score (Higher = Urgent)"
        ).add_to(m)
    
        folium.GeoJson(
            up_geojson,
            tooltip=folium.GeoJsonTooltip(
                fields=["district_name"],
                aliases=["District:"],
                sticky=True
            )
        ).add_to(m)
    
        st_folium(m, use_container_width=True, height=720)
        
                # ====================================================
        # 📋 DISTRICT PRIORITY LIST (HIGH → LOW)
        # ====================================================
        st.markdown("## 📋 District Inspection Priority — High to Low")
        
        priority_tbl = (
            district_df
            .sort_values("Inspection Score", ascending=False)
            .reset_index(drop=True)
        )
        
        priority_tbl.insert(0, "S. No.", range(1, len(priority_tbl) + 1))
        
        st.dataframe(
            priority_tbl,
            use_container_width=True,
            height=min(520, max(300, 32 * len(priority_tbl))),
            hide_index=True
        )
        
        st.download_button(
            "📤 Download District Priority List",
            priority_tbl.to_csv(index=False).encode(),
            file_name="district_inspection_priority_high_to_low.csv",
            mime="text/csv"
        )


    
        # ====================================================
        # 🔽 DISTRICT SELECTION (AUTHORITATIVE)
        # ====================================================
        st.markdown("## 🔽 Select District for Detailed Review")
    
        sel_dist = st.selectbox(
            "Select District",
            district_df["District"].unique()
        )
    
        sel_row = district_df[district_df["District"] == sel_dist].iloc[0]
    
        # ====================================================
        # 📊 INSPECTION SCORE — SELECTED DISTRICT (EXPLAINED)
        # ====================================================
        # ====================================================
        # 📊 INSPECTION SCORE — SELECTED DISTRICT (EXPLAINED)
        # ====================================================
        st.markdown("## 📊 Inspection Score — Selected District")
        
        score_breakdown = pd.DataFrame([
            ["Bottom 25% Institutes", sel_row["Bottom 25%"], 4, sel_row["Bottom 25%"] * 4],
            ["25–50% Performance Band", sel_row["25–50%"], 3, sel_row["25–50%"] * 3],
            ["50–75% Performance Band", sel_row["50–75%"], 2, sel_row["50–75%"] * 2],
            ["FSR Deficit (>1:25)", sel_row["FSR Deficit"], 2, sel_row["FSR Deficit"] * 2],
            ["Rank Jump (±50)", sel_row["Rank Jump ±50"], 3, sel_row["Rank Jump ±50"] * 3],
        ], columns=["Component", "Count", "Points / Unit", "Total Points"])
        
        # Add correct serial number
        score_breakdown.insert(0, "S. No.", range(1, len(score_breakdown) + 1))
        
        st.dataframe(
            score_breakdown,
            use_container_width=True,
            height=260,
            hide_index=True
        )

    
        final_score = sel_row["Inspection Score"]
    
        if final_score >= 15:
            status = "🔴 Physical Inspection REQUIRED"
            color = "#C62828"
            remark = "High concentration of poor-performing institutes and multiple risk indicators."
        elif final_score >= 8:
            status = "🟠 Inspection Recommended"
            color = "#EF6C00"
            remark = "Several warning indicators observed. Scheduled inspection advised."
        elif final_score >= 1:
            status = "🟡 Desk Review Sufficient"
            color = "#F9A825"
            remark = "Limited issues detected. Monitoring through records is sufficient."
        else:
            status = "🟢 Normal — No Inspection Needed"
            color = "#2E7D32"
            remark = "District performance is within acceptable limits."
    
        st.markdown(
            f"""
            <div style="padding:18px;border-radius:14px;
            background:#FFF8E1;border-left:8px solid {color};">
            <b>District:</b> {sel_dist}<br><br>
            <b>Final Inspection Score:</b> {final_score}<br><br>
    
            <b>Score Interpretation:</b><br>
            • 15+ → High Priority<br>
            • 8–14 → Medium Priority<br>
            • 1–7 → Low Priority<br>
            • 0 → Normal<br><br>
    
            <b>Decision:</b> {status}<br><br>
            <b>Comment:</b> {remark}
            </div>
            """,
            unsafe_allow_html=True
        )
    
           # ====================================================
        # 🏫 INSTITUTE LIST — SELECTED DISTRICT ONLY
        # ====================================================
        st.markdown("## 🏫 Institutes in Selected District")
        
        inst_tbl = (
            work[work["District_norm"] == sel_dist]
            [[COL_INST, COL_TOTAL, COL_TLR, COL_PPRPI, COL_DO, FSR_COL, COL_RANK]]
            .sort_values(COL_TOTAL)
            .reset_index(drop=True)
        )
        
        # Add proper serial number
        inst_tbl.insert(0, "S. No.", range(1, len(inst_tbl) + 1))
        
        st.dataframe(
            inst_tbl,
            use_container_width=True,
            height=min(520, max(260, 35 * len(inst_tbl))),
            hide_index=True
        )
        
        st.download_button(
            "📤 Download Institute List",
            inst_tbl.to_csv(index=False).encode(),
            file_name=f"{sel_dist}_inspection_institutes.csv",
            mime="text/csv"
        )

    
            
                # ====================================================
        # 📘 INSPECTION SCORE — OFFICIAL CALCULATION SCRIPT
        # ====================================================
        
        st.markdown("## 📘 Inspection Score — Calculation & Decision Logic")
        
        st.markdown(
            """
        ### 🎯 Objective
        To identify **districts requiring physical inspection** using only
        **verifiable institutional performance data filled by Institutes**.
        
        
        
        ---
        
        ### 🔢 Step 1: Performance Band Classification (Institute Level)
        
        Each institute is classified using **state-wide percentiles**:
        
        | Performance Band | Condition | Points per Institute |
        |------------------|----------|---------------------|
        | Bottom 25% | Score ≤ 25th percentile | **4** |
        | 25–50% | 25th < Score ≤ 50th | **3** |
        | 50–75% | 50th < Score ≤ 75th | **2** |
        | 75–100% | Score > 75th percentile | **0** |
        
        ---
        
        ### ⚠️ Step 2: Risk Indicators (Additional Points)
        
        | Risk Indicator | Condition | Points |
        |---------------|----------|--------|
        | Faculty–Student Ratio Deficit | Worse than **1:25** | **+2 per institute** |
        | Abrupt Rank Change | ±50 rank movement | **+3 per institute** |
        
        ---
        
        ### 🧮 Step 3: District Inspection Score Formula
        
                    Inspection Score = (4 × Bottom 25% Institutes)
                                     
                                       (3 × 25–50% Institutes)
        
                                       (2 × 50–75% Institutes)
        
                                       (2 × FSR Deficit Institutes)
        
                                       (3 × Rank Jump ±50 Institutes)
                        
                        
        ---
        )
        ### 🚦 Step 4: Inspection Decision Thresholds
        
        | Score Range | Priority Level | Action |
        |------------|---------------|--------|
        | **15 or more** | 🔴 High | Immediate physical inspection |
        | **8–14** | 🟠 Medium | Schedule inspection |
        | **1–7** | 🟡 Low | Desk review sufficient |
        | **0** | 🟢 Normal | No inspection required |
        
        ---
        
        ### 📌 Operational Notes
        • Score is calculated **district-wise**  
        • Only **participating institutes** are included  
        • All counts are **fact-based and filled by Institute**  
      
        
        
        """
        )
    # ========================================================
    # XAI (FEATURE IMPORTANCE)
    # ========================================================
      # ========================================================
        # TAB 8 — EXPLAINABLE AI (XAI) + AI FINDINGS + SUPPORT PLAN
        # ========================================================
          # ========================================================
    # TAB 8 — EXPLAINABLE AI (XAI) + CORRECT GAP-TO-100 ANALYSIS
    # ========================================================
    with tabs[8]:
        st.markdown("## 🤖 Explainable AI (XAI): Assessment Findings and Improvement Opportunities")
    
        st.caption(
            """
            This section explains **why an institute has not achieved the full 100 marks**,
            using **Explainable AI** and the **actual evaluation framework**.
    
            ✔ What matters most at state level  
            ✔ What is reducing marks for the selected institute  
            ✔ Exactly **how many marks are lost** from each parameter  
            ✔ What action will give **maximum score improvement**
            """
        )
    
        # ====================================================
        # SAFE SHAP IMPORT
        # ====================================================
        try:
            import shap
            SHAP_AVAILABLE = True
        except Exception:
            SHAP_AVAILABLE = False
    
        # ====================================================
        # DATA PREPARATION
        # ====================================================
        work = df_m[df_m[COL_TOTAL] > 0].copy()
        X = work[[COL_TLR, COL_PPRPI, COL_DO]]
        y = work[COL_TOTAL]
    
        # ====================================================
        # MODEL TRAINING
        # ====================================================
        model = RandomForestRegressor(
            n_estimators=300,
            random_state=42,
            max_depth=6
        )
        model.fit(X, y)
    
        # ====================================================
        # 🌍 GLOBAL XAI — STATE LEVEL
        # ====================================================
        st.markdown("## 🌍 State-Level Insight: What Drives Scores")
    
        if SHAP_AVAILABLE:
            explainer = shap.TreeExplainer(model)
            shap_values = explainer.shap_values(X)
    
            global_imp = pd.DataFrame({
                "Parameter": X.columns,
                "Influence (SHAP)": np.abs(shap_values).mean(axis=0)
            })
        else:
            global_imp = pd.DataFrame({
                "Parameter": X.columns,
                "Influence (Model)": model.feature_importances_
            })
    
        global_imp = global_imp.sort_values(
            global_imp.columns[1], ascending=False
        ).reset_index(drop=True)
    
        fig_global = px.bar(
            global_imp,
            x="Parameter",
            y=global_imp.columns[1],
            color=global_imp.columns[1],
            color_continuous_scale="YlOrRd",
            title="State Wise Importance — Parameters Influencing Total Score"
        )
        st.plotly_chart(fig_global, use_container_width=True)
    
        st.info(
            f"""
            **Explainable AI Finding (State Level)**  
            • **{global_imp.iloc[0]['Parameter']}** has the strongest impact statewide  
            • Improving this parameter gives **maximum return on effort**
            """
        )
    
        # ====================================================
        # 🎯 INSTITUTE SELECTION
        # ====================================================
        st.markdown("## 🎯 Institute-Level Analysis")
    
        sel_inst = st.selectbox(
            "Select Institute",
            sorted(work[COL_INST].unique())
        )
    
        inst_row = work[work[COL_INST] == sel_inst]
        inst_X = inst_row[[COL_TLR, COL_PPRPI, COL_DO]]
    
        actual_total = inst_row[COL_TOTAL].values[0]
        gap_to_100 = 100 - actual_total
    
        # ====================================================
        # LOCAL XAI (WHY SCORE IS REDUCED)
        # ====================================================
        if SHAP_AVAILABLE:
            inst_shap = explainer.shap_values(inst_X)[0]
        else:
            inst_shap = (inst_X.iloc[0] - X.mean()).values
    
        local_xai = pd.DataFrame({
            "Parameter": inst_X.columns,
            "Impact on Total Score": inst_shap
        }).sort_values("Impact on Total Score")
    
        fig_local = px.bar(
            local_xai,
            x="Impact on Total Score",
            y="Parameter",
            orientation="h",
            color="Impact on Total Score",
            color_continuous_scale="RdYlGn",
            title="Performance Drivers & Gaps at Institute Level"
        )
        st.plotly_chart(fig_local, use_container_width=True)
    
        # ====================================================
        # ✅ CORRECT GAP-TO-100 (WEIGHTED)
            # ====================================================
        st.markdown("## 📉 Institute Assessment Review: Observations and Improvement Areas")
        
        WEIGHTS = {
            COL_TLR: 30,
            COL_PPRPI: 35,
            COL_DO: 35
        }
        
        gap_rows = []
        for param, max_marks in WEIGHTS.items():
            actual = inst_row[param].values[0]
            lost = max_marks - actual
        
            gap_rows.append({
                "Parameter": param,
                "Maximum Possible Marks": max_marks,
                "Marks Achieved": actual,
                "Marks Lost": lost
            })
        
        gap_df = pd.DataFrame(gap_rows)
        
        fig_gap = px.bar(
            gap_df,
            x="Parameter",
            y=["Marks Lost", "Marks Achieved"],   # 🔥 FIX ORDER (NO OVERRIDE)
            barmode="group",
            title="Marks Achieved vs Marks Lost",
            color_discrete_sequence=[
                "#60A5FA",  # 🔵 Marks Lost (Light Blue)
                "#F59E0B"   # 🟠 Marks Achieved (Muted Amber)
            ]
        )
        
        fig_gap.update_traces(
            texttemplate="%{y:.2f}",
            textposition="outside",
            cliponaxis=False
        )
        
        fig_gap.update_layout(
            yaxis_title="Marks",
            xaxis_title="Assessment Parameter",
            plot_bgcolor="white",
            paper_bgcolor="white",
        
            # ✅ LEGEND IN TWO ROWS (FORCED)
            legend=dict(
                orientation="h",
                yanchor="top",
                y=-0.42,
                xanchor="center",
                x=0.5,
                title_text="Score Components",
                font=dict(size=14),
        
                entrywidth=140,              # 🔥 KEY: forces wrap
                entrywidthmode="pixels"
            ),
        
            # ✅ EXTRA SPACE FOR 2 ROW LEGEND
            margin=dict(
                t=80,
                b=200,                       # 🔥 IMPORTANT
                l=50,
                r=40
            ),
        
            bargap=0.25,
            bargroupgap=0.15
        )

        
        st.plotly_chart(fig_gap, use_container_width=True)



    
        # ====================================================
        # 📊 GAP EXPLANATION TABLE
        # ====================================================
        st.markdown("## 📊 Parameter-Wise Assessment Breakdown")
    
        st.dataframe(
            gap_df.sort_values("Marks Lost", ascending=False),
            use_container_width=True,
            height=220
        )
    
        # ====================================================
        # 🛠️ AUTO ACTION PLAN — WHERE TO WORK MORE
        # ====================================================
        st.markdown("## 🛠️ Priority Areas for Institutional Development")
        worst = gap_df.sort_values("Marks Lost", ascending=False).iloc[0]
        
        st.error(
            f"""
            📌 **Primary Gap Identified**
        
            • Maximum marks lost in **{worst['Parameter']}**  
            • Lost **{worst['Marks Lost']:.2f} marks** out of {worst['Maximum Possible Marks']:.2f}
        
            👉 Focusing improvement here will give the **highest score jump**
            """
        )

    
        # ====================================================
        # 🎓 FINAL OFFICER-FRIENDLY SUMMARY
        # ====================================================
        st.info(
            """
            🧠 **How to Read This Section**
    
            • State Wise XAI → what matters most for policy  
            • Institute Wise XAI → why this institute is weak  
            • Gap-to-100 → exact marks missing  
            • Action focus → where effort gives maximum gain  
    
            🎯 Objective: **Improve institutes, not penalize them**
            """
        )




    # ========================================================
    # EXPORT
    # ========================================================
    with tabs[9]:
        st.markdown("## 📤 Export Reports")
    
        if st.button("📊 Download PPT Report"):
            ppt = generate_sirf_ppt(
                session_label=selected_month.strftime("%b-%Y"),
                kpis={
                    "total": total_institutes_master,
                    "participated": participated,
                    "not_participated": not_participated,
                    "avg_score": avg_score
                },
                fig_grade=fig_grade,
                fig_zone=figz,
                participation_df=summary.head(20),  # or full
                red_flags_df=red_flags if 'red_flags' in locals() else pd.DataFrame()
            )
    
            st.download_button(
                "⬇️ Download SIRF PPT",
                ppt,
                file_name=f"SIRF_Report_{selected_month.strftime('%b_%Y')}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
    

# ============================================================
if __name__ == "__main__":
    main()


# In[ ]:





# In[ ]:





# In[ ]:




